import fitz
import docx
from pptx import Presentation


def extract_pdf_text(file_path):
    text = ""

    pdf = fitz.open(file_path)

    for page in pdf:
        text += page.get_text()

    return text


def extract_docx_text(file_path):
    document = docx.Document(file_path)

    text = "\n".join(
        paragraph.text
        for paragraph in document.paragraphs
    )

    return text


def extract_pptx_text(file_path):
    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def extract_text(file_path, extension):
    if extension == ".pdf":
        return extract_pdf_text(file_path)

    elif extension == ".docx":
        return extract_docx_text(file_path)

    elif extension == ".pptx":
        return extract_pptx_text(file_path)

    else:
        raise ValueError("Unsupported file type")